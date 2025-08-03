import csv
import html
from io import BytesIO, StringIO
from typing import Callable, Tuple

from bs4 import BeautifulSoup
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table as DocxTable
from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
import xlrd
from xlrd.sheet import Sheet

from libs.logging import get_logger
from process_document.types import Table

from .blob_storage import BlobFile

logger = get_logger(__name__)
logger.setLevel("INFO")


def process_pdf_files_for_ursa(
    file_data: bytes, blob_path: str, parser_func: Callable[[BytesIO], str]
) -> Tuple[list[BlobFile], int]:
    """
    PDFファイルを処理します。

    :param file_data: fileデータ
    :type file_data: bytes
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: PDFファイルのリスト
    :rtype: list[BlobFile]
    """

    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        pdf_reader = PdfReader(bytes_stream)

        combined_text = ""
        for page in pdf_reader.pages:
            pdf_writer = PdfWriter()
            with BytesIO() as page_bytes_stream:
                pdf_writer.add_page(page)
                pdf_writer.write(page_bytes_stream)
                page_bytes_stream.seek(0)
                combined_text += parser_func(page_bytes_stream)
        uploaded_files.append({"file_name": blob_path, "text": combined_text, "page_number": 0})

        return uploaded_files, len(pdf_reader.pages)


def _find_table_from_element(tables_in_docx: list[DocxTable], element: CT_Tbl) -> DocxTable | None:
    for table in tables_in_docx:
        if table._tbl == element:
            return table
    return None  # テーブルが見つからない場合


def convert_docx_table_to_html(table: DocxTable) -> str:
    try:
        soup = BeautifulSoup("", "html.parser")
        table_soup = soup.new_tag("table")
        for row in table.rows:
            tr_tag = soup.new_tag("tr")
            table_soup.append(tr_tag)
            for cell in row.cells:
                new_tag = soup.new_tag("td")
                new_tag.append(html.escape(cell.text or ""))
                tr_tag.append(new_tag)
            table_soup.append(tr_tag)
        return table_soup.prettify() + "\n"
    finally:
        soup.decompose()


def process_word_files_for_ursa(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    docxファイルを処理します。

    :param file_data: fileデータ
    :type file_data: str
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: docxファイルのリスト
    :rtype: list[BlobFile]
    """

    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        doc = Document(bytes_stream)
        tables_in_docx = doc.tables if doc.tables else []

        text = ""
        for element in doc.element.body:
            # elementの型はCT_PかCT_TblかCT_SectPrのいずれか
            # CT_P: パラグラフ（文章）
            # CT_Tbl: テーブル（表）
            # CT_SectPr: レイアウト設定(用紙サイズ、ページの余白等)
            if isinstance(element, CT_P):
                text += element.text + "\n"
                continue
            if isinstance(element, CT_Tbl):
                table_docx = _find_table_from_element(tables_in_docx, element)
                if not table_docx:
                    continue
                try:
                    table_text = Table.from_docx_table(table_docx).to_text() + "\n"
                except Exception as e:
                    logger.warning(f"Failed to convert docx table to table class: {table_docx}", exc_info=e)
                    table_text = convert_docx_table_to_html(table_docx)
                text += table_text
        uploaded_files.append({"file_name": blob_path, "text": text, "page_number": 0})
        return uploaded_files


def convert_xlsx_table_to_csv(table: Worksheet) -> str:
    existing_max_row, existing_max_col = 0, 0
    for _row in range(table.max_row, 0, -1):
        if any(
            table.cell(row=_row, column=_col).value is not None
            and str(table.cell(row=_row, column=_col).value).strip()
            for _col in range(1, table.max_column + 1)
        ):
            existing_max_row = _row
            break
    for _col in range(table.max_column, 0, -1):
        if any(
            table.cell(row=row, column=_col).value is not None and str(table.cell(row=_row, column=_col).value).strip()
            for row in range(1, table.max_row + 1)
        ):
            existing_max_col = _col
            break
    try:
        output = StringIO()
        csv_writer = csv.writer(output)

        for row in table.iter_rows(max_row=existing_max_row, max_col=existing_max_col):
            csv_row = [str(cell.value) if cell.value is not None else "" for cell in row]
            csv_writer.writerow(csv_row)
        return "<table>\n" + output.getvalue() + "\n</table>\n"
    finally:
        output.close()


def convert_xls_table_to_csv(sheet: Sheet) -> str:
    existing_max_row, existing_max_col = 0, 0
    for _row in range(sheet.nrows - 1, -1, -1):
        if any(sheet.cell_value(_row, _col) not in (None, "") for _col in range(sheet.ncols)):
            existing_max_row = _row + 1
            break
    for _col in range(sheet.ncols - 1, -1, -1):
        if any(sheet.cell_value(_row, _col) not in (None, "") for _row in range(sheet.nrows)):
            existing_max_col = _col + 1
            break
    try:
        output = StringIO()
        csv_writer = csv.writer(output)

        for row_idx in range(existing_max_row):
            csv_row = [
                str(sheet.cell_value(row_idx, col_idx)) if sheet.cell_value(row_idx, col_idx) not in (None, "") else ""
                for col_idx in range(existing_max_col)
            ]
            csv_writer.writerow(csv_row)

        return "<table>\n" + output.getvalue() + "\n</table>\n"
    finally:
        output.close()


def process_excel_files_for_ursa(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    xlsxファイルを処理します。

    :param file_data: fileデータ
    :type file_data: bytes
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: excelファイルのリスト
    :rtype: list[BlobFile]
    """

    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        workbook = None
        try:
            workbook = load_workbook(bytes_stream, data_only=True)

            text = ""
            for _i, sheet_name in enumerate(workbook.sheetnames):
                worksheet = workbook[sheet_name]
                try:
                    table_text = Table.from_xlsx_table(worksheet, sheet_name).to_text() + "\n"
                except Exception as e:
                    logger.warning(f"Failed to convert xlsx table to table class: {worksheet}", exc_info=e)
                    table_text = convert_xlsx_table_to_csv(worksheet)
                text += table_text
            uploaded_files.append({"file_name": blob_path, "text": text, "page_number": 0})
            return uploaded_files
        finally:
            if workbook:
                workbook.close()


def process_excel_xls_files_for_ursa(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    xlsファイルを処理します。

    :param file_data: fileデータ
    :type file_data: bytes
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: excelファイルのリスト
    :rtype: list[BlobFile]
    """
    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        workbook = xlrd.open_workbook(file_contents=bytes_stream.read())

        text = ""
        for sheet_index in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_index)
            sheet_name = workbook.sheet_names()[sheet_index]
            try:
                table_text = Table.from_xls_table(sheet, sheet_name).to_text() + "\n"
            except Exception as e:
                logger.warning(f"Failed to convert xls table to table class: {sheet_name}", exc_info=e)
                table_text = convert_xls_table_to_csv(sheet)
            text += table_text
        uploaded_files.append({"file_name": blob_path, "text": text, "page_number": 0})
    return uploaded_files


def process_powerpoint_files_for_ursa(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    txtファイルを取得します。
    すべてのスライドのテキストを1つのBlobFileにまとめて返します。
    :param file_data: fileデータ
    :type file_data: bytes
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: pptxファイルのリスト
    :rtype: list[BlobFile]
    """

    all_slides_text = ""
    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        prs = Presentation(bytes_stream)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_slides_text += shape.text + "\n"
        uploaded_files.append({"file_name": blob_path, "text": all_slides_text, "page_number": 0})
        return uploaded_files
