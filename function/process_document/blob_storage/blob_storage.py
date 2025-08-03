from concurrent.futures import ThreadPoolExecutor, as_completed
import csv
import html
from io import BytesIO, StringIO
from typing import Callable, Tuple, TypedDict

from bs4 import BeautifulSoup
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table as DocxTable
from langchain.text_splitter import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
import xlrd
from xlrd.sheet import Sheet

from libs.logging import get_logger
from process_document.document_intelligence.document_intelligence import parse_pdf_by_document_intelligence
from process_document.libs.document import decode_file
from process_document.llm_document_reader.llm_document_reader import parse_pdf_by_llm_document_reader
from process_document.types import DEFAULT_TABLE_CHUNK_SIZE, Table

from .client import get_blob_container

logger = get_logger(__name__)

TXT_CHUNK_SIZE = 500
TXT_CHUNK_OVERLAP = 50
DEFAULT_ENCODING_NAME = "o200k_base"


class BlobFile(TypedDict):
    file_name: str
    page_number: int
    text: str


def get_files_from_blob(container_name: str, blob_path: str) -> bytes:
    """
    blob_strageからファイルを取得します。

    :param container_name: コンテナ名
    :type container_name: str
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: ファイルを読み込んだデータ
    :rtype: bytes
    """
    blob_container = get_blob_container(container_name)
    blob_client = blob_container.get_blob_client(blob_path)
    with BytesIO() as bytes_stream:
        blob_client.download_blob().readinto(bytes_stream)
        bytes_stream.seek(0)
        return bytes_stream.read()


def process_pdf_files(
    file_data: bytes, blob_path: str, parser_func: Callable[[BytesIO], str]
) -> Tuple[list[BlobFile], int]:
    """
    PDFファイルを処理します。

    :param file_data: fileデータ
    :type file_data: bytes
    :param blob_path: Blobのパス
    :type blob_path: str
    :param parser_func: パーサー関数
    :type parser_func: Callable[[BytesIO], str]
    :return: PDFファイルのリストとページ数
    :rtype: Tuple[list[BlobFile], int]
    """
    with BytesIO(file_data) as bytes_stream:
        pdf_reader = PdfReader(bytes_stream)
        num_pages = len(pdf_reader.pages)

        page_streams: list[bytes] = []
        for page_number in range(num_pages):
            pdf_writer = PdfWriter()
            pdf_writer.append(pdf_reader, pages=(page_number, page_number + 1))
            with BytesIO() as page_bytes_stream:
                pdf_writer.write(page_bytes_stream)
                page_bytes_stream.seek(0)
                page_streams.append(page_bytes_stream.read())

        def process_page(page_number: int) -> BlobFile:
            with BytesIO(page_streams[page_number]) as page_bytes_stream:
                text = parser_func(page_bytes_stream)
                return {"file_name": blob_path, "page_number": page_number + 1, "text": text}

        MAX_WORKERS = 5
        uploaded_files = []
        if parser_func in {parse_pdf_by_document_intelligence, parse_pdf_by_llm_document_reader}:
            with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                futures = [executor.submit(process_page, i) for i in range(num_pages)]
                for future in as_completed(futures):
                    uploaded_files.append(future.result())
        else:
            for i in range(num_pages):
                uploaded_files.append(process_page(i))

        return uploaded_files, num_pages


def _find_table_from_element(tables_in_docx: list[DocxTable], element: CT_Tbl) -> DocxTable | None:
    for table in tables_in_docx:
        if table._tbl == element:
            return table
    return None  # テーブルが見つからない場合


def convert_docx_table_to_html(table: DocxTable) -> str:
    try:
        soup = BeautifulSoup("", "html.parser")
        table_soup = soup.new_tag("table")
        for row in table.rows:
            tr_tag = soup.new_tag("tr")
            table_soup.append(tr_tag)
            for cell in row.cells:
                new_tag = soup.new_tag("td")
                new_tag.append(html.escape(cell.text or ""))
                tr_tag.append(new_tag)
            table_soup.append(tr_tag)
        return table_soup.prettify() + "\n"
    finally:
        soup.decompose()


def process_word_files(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    docxファイルを処理します。

    :param file_data: fileデータ
    :type file_data: str
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: docxファイルのリスト
    :rtype: list[BlobFile]
    """

    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        doc = Document(bytes_stream)
        tables_in_docx = doc.tables if doc.tables else []

        text = ""
        for element in doc.element.body:
            # elementの型はCT_PかCT_TblかCT_SectPrのいずれか
            # CT_P: パラグラフ（文章）
            # CT_Tbl: テーブル（表）
            # CT_SectPr: レイアウト設定(用紙サイズ、ページの余白等)
            if isinstance(element, CT_P):
                text += element.text + "\n"
                continue
            if isinstance(element, CT_Tbl):
                table_docx = _find_table_from_element(tables_in_docx, element)
                if not table_docx:
                    continue
                try:
                    table_text = "\n".join(Table.from_docx_table(table_docx).to_chunks())
                except Exception as e:
                    logger.warning(f"Failed to convert docx table to table class: {table_docx}", exc_info=e)
                    table_text = convert_docx_table_to_html(table_docx)
                text += table_text
        uploaded_files.append({"file_name": blob_path, "text": text, "page_number": 0})
        return uploaded_files


def convert_xlsx_table_to_csv(table: Worksheet, sheet_name: str) -> list[str]:
    existing_max_row, existing_max_col = 0, 0
    for _row in range(table.max_row, 0, -1):
        if any(
            table.cell(row=_row, column=_col).value is not None
            and str(table.cell(row=_row, column=_col).value).strip()
            for _col in range(1, table.max_column + 1)
        ):
            existing_max_row = _row
            break
    for _col in range(table.max_column, 0, -1):
        if any(
            table.cell(row=row, column=_col).value is not None and str(table.cell(row=_row, column=_col).value).strip()
            for row in range(1, table.max_row + 1)
        ):
            existing_max_col = _col
            break
    try:
        output = StringIO()
        csv_writer = csv.writer(output)

        for row in table.iter_rows(max_row=existing_max_row, max_col=existing_max_col):
            csv_row = [str(cell.value) if cell.value is not None else "" for cell in row]
            csv_writer.writerow(csv_row)
        table_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            chunk_size=DEFAULT_TABLE_CHUNK_SIZE,
            chunk_overlap=0,
            separators=["\n\n", "\n", " ", ""],
            disallowed_special=(),
        )
        table_texts = table_splitter.split_text(output.getvalue())
        chunks = [f"<table>\n{table_text}\nTable: {sheet_name}\n</table>\n" for table_text in table_texts]
        return chunks
    finally:
        output.close()


def convert_xls_table_to_csv(sheet: Sheet, sheet_name: str) -> list[str]:
    existing_max_row, existing_max_col = 0, 0
    for _row in range(sheet.nrows - 1, -1, -1):
        if any(sheet.cell_value(_row, _col) not in (None, "") for _col in range(sheet.ncols)):
            existing_max_row = _row + 1
            break
    for _col in range(sheet.ncols - 1, -1, -1):
        if any(sheet.cell_value(_row, _col) not in (None, "") for _row in range(sheet.nrows)):
            existing_max_col = _col + 1
            break
    try:
        output = StringIO()
        csv_writer = csv.writer(output)
        for row_idx in range(existing_max_row):
            csv_row = [
                str(sheet.cell_value(row_idx, col_idx)) if sheet.cell_value(row_idx, col_idx) not in (None, "") else ""
                for col_idx in range(existing_max_col)
            ]
            csv_writer.writerow(csv_row)
        table_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            chunk_size=DEFAULT_TABLE_CHUNK_SIZE,
            chunk_overlap=0,
            separators=["\n\n", "\n", " ", ""],
            disallowed_special=(),
        )
        table_texts = table_splitter.split_text(output.getvalue())

        chunks = [f"<table>\n{table_text}\nTable: {sheet_name}\n</table>\n" for table_text in table_texts]
        return chunks
    finally:
        output.close()


def process_excel_files(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    xlsxファイルを処理します。

    :param file_data: fileデータ
    :type file_data: bytes
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: excelファイルのリスト
    :rtype: list[BlobFile]
    """

    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        workbook = None
        try:
            workbook = load_workbook(bytes_stream, data_only=True)

            for i, sheet_name in enumerate(workbook.sheetnames):
                worksheet = workbook[sheet_name]
                try:
                    table = Table.from_xlsx_table(worksheet, sheet_name)
                    chunks = table.to_chunks()
                    for chunk in chunks:
                        uploaded_files.append({"file_name": blob_path, "text": chunk, "page_number": i + 1})
                except Exception as e:
                    logger.warning(f"Failed to convert xlsx table to table class: {worksheet}", exc_info=e)
                    chunks = convert_xlsx_table_to_csv(worksheet, sheet_name)
                    for chunk in chunks:
                        uploaded_files.append({"file_name": blob_path, "text": chunk, "page_number": i + 1})
            return uploaded_files
        finally:
            if workbook:
                workbook.close()


def process_excel_xls_files(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    xlsファイルを処理します。

    :param file_data: fileデータ
    :type file_data: bytes
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: excelファイルのリスト
    :rtype: list[BlobFile]
    """
    uploaded_files: list[BlobFile] = []
    with BytesIO(file_data) as bytes_stream:
        workbook = xlrd.open_workbook(file_contents=bytes_stream.read())

        for sheet_index in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_index)
            sheet_name = workbook.sheet_names()[sheet_index]
            try:
                # Tableクラスで処理
                chunks = Table.from_xls_table(sheet, sheet_name).to_chunks()
            except Exception as e:
                logger.warning(f"Failed to convert xls table to table class: {sheet_name}", exc_info=e)
                # CSVに変換
                chunks = convert_xls_table_to_csv(sheet, sheet_name)
            for chunk in chunks:
                uploaded_files.append({"file_name": blob_path, "text": chunk, "page_number": sheet_index + 1})
    return uploaded_files


def process_powerpoint_files(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    :param file_data: fileデータ
    :type file_data: bytes
    :type container_name: str
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: pptxファイルのリスト
    :rtype: list[BlobFile]
    """
    uploaded_files: list[BlobFile] = []
    text_in_slide = ""
    with BytesIO(file_data) as bytes_stream:
        prs = Presentation(bytes_stream)
        # 各スライドをループ処理
        for slide_number, slide in enumerate(prs.slides):
            # スライド内のシェイプ（テキストボックスなど）をループ処理
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # テキストフレームを持つシェイプのテキストを取得
                    text_in_slide += shape.text
            uploaded_files.append({"file_name": blob_path, "text": text_in_slide, "page_number": slide_number + 1})
            text_in_slide = ""
        return uploaded_files


def process_text_files(file_data: bytes, blob_path: str) -> list[BlobFile]:
    """
    txtファイルを取得します。
    :param file_data: fileデータ
    :type file_data: BytesIO
    :param blob_path: Blobのパス
    :type blob_path: str
    :return: txtファイルのリスト
    :rtype: list[BlobFile]
    """

    text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name=DEFAULT_ENCODING_NAME,
        chunk_size=TXT_CHUNK_SIZE,
        chunk_overlap=TXT_CHUNK_OVERLAP,
        disallowed_special=(),
    )

    uploaded_files: list[BlobFile] = []
    text = decode_file(file_data)
    chunk = text_splitter.split_text(text)
    for i, chunked_text in enumerate(chunk):
        uploaded_files.append({"file_name": blob_path, "text": chunked_text, "page_number": i + 1})
    return uploaded_files
