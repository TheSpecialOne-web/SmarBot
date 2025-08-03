import binascii
import csv
import html
from io import BytesIO, StringIO
from zipfile import BadZipFile

from bs4 import BeautifulSoup
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table as DocxTable
from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from PIL import Image
from pptx import Presentation
from pydantic import RootModel, StrictInt, StrictStr
from pypdf import PdfReader
from pypdf.errors import FileNotDecryptedError, PdfReadError

from api.libs.exceptions import BadRequest
from api.libs.logging import get_logger
from api.libs.table import Table

logger = get_logger()


class Content(RootModel):
    root: StrictStr

    def merge(self, to_merge: "Content") -> "Content":
        return Content(root=self.root + to_merge.root)


class ContentPageCount(RootModel):
    root: StrictInt


class BlobContent(RootModel):
    root: bytes

    def count_pages(self) -> ContentPageCount:
        with BytesIO(self.root) as pdf_stream:
            try:
                reader = PdfReader(pdf_stream)
            except PdfReadError:
                raise BadRequest("ファイルが破損している可能性があります。")
            except BadZipFile:
                raise BadRequest("ファイルが破損している可能性があります。")
            try:
                return ContentPageCount(root=len(reader.pages))
            except FileNotDecryptedError:
                raise BadRequest("PDFファイルが暗号化されているため、読み込めません。")
            except binascii.Error:
                raise BadRequest("ファイルが破損している可能性があります。")

    def parse_pdf_file_by_pypdf(self) -> Content:
        with BytesIO(self.root) as pdf_stream:
            try:
                reader = PdfReader(pdf_stream)
            except PdfReadError:
                raise BadRequest("ファイルが破損している可能性があります。")
            except BadZipFile:
                raise BadRequest("ファイルが破損している可能性があります。")
            blob_content = ""
            try:
                for page in range(len(reader.pages)):
                    blob_content += reader.pages[page].extract_text()
            except FileNotDecryptedError:
                raise BadRequest("PDFファイルが暗号化されているため、読み込めません。")
            except binascii.Error:
                raise BadRequest("ファイルが破損している可能性があります。")
            return Content(root=blob_content)

    def _convert_docx_table_to_html(self, table: DocxTable) -> str:
        try:
            soup = BeautifulSoup("", "html.parser")
            table_soup = soup.new_tag("table")
            for row in table.rows:
                tr_tag = soup.new_tag("tr")
                table_soup.append(tr_tag)
                for cell in row.cells:
                    new_tag = soup.new_tag("td")
                    new_tag.append(html.escape(cell.text or ""))
                    tr_tag.append(new_tag)
                table_soup.append(tr_tag)
            return table_soup.prettify() + "\n"
        finally:
            soup.decompose()

    @staticmethod
    def _find_table_from_element(tables_in_docx: list[DocxTable], element: CT_Tbl) -> DocxTable | None:
        for table in tables_in_docx:
            if table._tbl == element:
                return table
        return None  # テーブルが見つからない場合

    def parse_docx_file(self) -> Content:
        with BytesIO(self.root) as stream:
            try:
                doc = Document(stream)
            except ValueError:
                raise BadRequest("ファイルの形式が正しくありません。")
            except BadZipFile:
                raise BadRequest("ファイルが破損している可能性があります。")

            tables_in_docx = doc.tables if doc.tables else []

            text = ""
            for element in doc.element.body:
                # elementの型はCT_PかCT_TblかCT_SectPrのいずれか
                # CT_P: パラグラフ（文章）
                # CT_Tbl: テーブル（表）
                # CT_SectPr: レイアウト設定(用紙サイズ、ページの余白等)
                if isinstance(element, CT_P):
                    text += element.text + "\n"
                    continue
                if isinstance(element, CT_Tbl):
                    table_docx = self._find_table_from_element(tables_in_docx, element)
                    if not table_docx:
                        continue
                    try:
                        table = Table.from_docx_table(table_docx)
                        table_text = table.to_text() + "\n"
                    except Exception as e:
                        logger.warning(f"Failed to convert docx table to table class: {table_docx}", exc_info=e)
                        table_text = self._convert_docx_table_to_html(table_docx)
                    text += table_text
            return Content(root=text)

    def _convert_xlsx_table_to_csv(self, table: Worksheet) -> str:
        existing_max_row, existing_max_col = 0, 0
        for _row in range(table.max_row, 0, -1):
            if any(
                table.cell(row=_row, column=_col).value is not None
                and str(table.cell(row=_row, column=_col).value).strip()
                for _col in range(1, table.max_column + 1)
            ):
                existing_max_row = _row
                break
        for _col in range(table.max_column, 0, -1):
            if any(
                table.cell(row=row, column=_col).value is not None
                and str(table.cell(row=_row, column=_col).value).strip()
                for row in range(1, table.max_row + 1)
            ):
                existing_max_col = _col
                break
        try:
            output = StringIO()
            csv_writer = csv.writer(output)

            for row in table.iter_rows(max_row=existing_max_row, max_col=existing_max_col):
                csv_row = [str(cell.value) if cell.value is not None else "" for cell in row]
                csv_writer.writerow(csv_row)
            return "<table>\n" + output.getvalue() + "\n</table>\n"
        finally:
            output.close()

    def parse_xlsx_file(self) -> Content:
        with BytesIO(self.root) as stream:
            workbook = None
            try:
                workbook = load_workbook(filename=stream, data_only=True)
                contents: list[str] = []
                for sheet_name in workbook.sheetnames:
                    worksheet = workbook[sheet_name]
                    try:
                        content = Table.from_xlsx_table(worksheet, sheet_name).to_text()
                    except Exception as e:
                        logger.warning(f"Failed to convert {sheet_name} to xlsx table class: {worksheet}", exc_info=e)
                        content = self._convert_xlsx_table_to_csv(worksheet)
                    contents.append(content)
                return Content(root="\n".join(contents))
            except ValueError:
                raise BadRequest("ファイルの形式が正しくありません。")
            except BadZipFile:
                raise BadRequest("ファイルが破損している可能性があります。")
            finally:
                if workbook:
                    workbook.close()

    def parse_pptx_file(self) -> Content:
        with BytesIO(self.root) as stream:
            try:
                prs = Presentation(stream)
            except ValueError:
                raise BadRequest("ファイルの形式が正しくありません。")
            except BadZipFile:
                raise BadRequest("ファイルが破損している可能性があります。")
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text
            return Content(root=content)

    def resize_image(self, min_size: tuple[int, int], max_size: tuple[int, int], format_: str = "jpeg"):
        with BytesIO(self.root) as bytes_stream, Image.open(bytes_stream) as img:
            width, height = img.size
            max_width, max_height = max_size
            min_width, min_height = min_size
            if min_width <= width <= max_width and min_height <= height <= max_height:
                return

            aspect_ratio = width / height
            # resize based on width
            new_width = min(max(width, min_width), max_width)
            new_height = int(new_width / aspect_ratio)
            # resize based on height
            if new_height < min_height or new_height > max_height:
                new_height = min(max(height, min_height), max_height)
                new_width = int(new_height * aspect_ratio)

            with img.resize((new_width, new_height), Image.Resampling.LANCZOS) as resized_img:
                with BytesIO() as output:
                    resized_img.save(output, format=format_)
                    self.root = output.getvalue()
